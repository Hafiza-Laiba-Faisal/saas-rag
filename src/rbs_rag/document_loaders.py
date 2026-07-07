from __future__ import annotations

import csv
import hashlib
import html.parser
import io
import re
import traceback
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Optional
from zipfile import ZipFile

from .models import LoadedDocument

SUPPORTED_EXTENSIONS = {".txt", ".md", ".markdown", ".html", ".htm", ".docx", ".pdf", ".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp", ".xlsx", ".xls", ".pptx", ".ppt", ".csv"}


def iter_document_files(path: Path) -> list[Path]:
    if path.is_file():
        return [path]
    files = [item for item in path.rglob("*") if item.is_file() and item.suffix.lower() in SUPPORTED_EXTENSIONS]
    return sorted(files)


def load_document(path: Path, metadata: dict | None = None, use_ocr: bool = False) -> LoadedDocument:
    suffix = path.suffix.lower()
    text = ""

    if suffix in {".txt", ".md", ".markdown"}:
        text = path.read_text(encoding="utf-8", errors="replace")
    elif suffix in {".html", ".htm"}:
        text = _strip_html(path.read_text(encoding="utf-8", errors="replace"))
    elif suffix == ".docx":
        text = _read_docx(path)
    elif suffix in {".pptx", ".ppt"}:
        text = _read_pptx(path)
    elif suffix == ".csv":
        text = _read_csv(path)
    elif suffix == ".pdf":
        text = _read_pdf(path, use_ocr=use_ocr)
    elif suffix == ".xlsx":
        text = _read_xlsx(path)
    elif suffix == ".xls":
        text = _read_xls(path)
    elif suffix in {".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp"}:
        text = _read_image_with_ocr(path)
    else:
        raise ValueError(f"Unsupported document type: {path.suffix}")

    base_metadata = {
        "document_name": path.name,
        "document_type": suffix.lstrip(".") or "text",
        "source_path": str(path),
        "ocr_applied": use_ocr or suffix in {".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp"},
    }
    if metadata:
        base_metadata.update(metadata)
    return LoadedDocument(
        document_id=_document_id(path),
        path=str(path),
        name=path.name,
        document_type=suffix.lstrip(".") or "text",
        text=text,
        metadata=base_metadata,
    )


def _document_id(path: Path) -> str:
    return hashlib.sha1(str(path.resolve()).encode("utf-8")).hexdigest()


def _strip_html(source: str) -> str:
    parser = _HTMLTextParser()
    parser.feed(source)
    return parser.text()


class _HTMLTextParser(html.parser.HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._parts: list[str] = []

    def handle_data(self, data: str) -> None:
        cleaned = data.strip()
        if cleaned:
            self._parts.append(cleaned)

    def text(self) -> str:
        return "\n".join(self._parts)


def _read_docx(path: Path) -> str:
    paragraphs: list[str] = []
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    with ZipFile(path) as archive:
        xml = archive.read("word/document.xml")
    root = ET.fromstring(xml)
    for paragraph in root.findall(".//w:p", namespace):
        texts = [node.text for node in paragraph.findall(".//w:t", namespace) if node.text]
        if texts:
            paragraphs.append("".join(texts))
    return "\n".join(paragraphs)


def _read_pptx(path: Path) -> str:
    """Extract text from all slides in a .pptx file."""
    from pptx import Presentation
    prs = Presentation(str(path))
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    line = "\t".join(cells)
                    if line.strip():
                        slide_texts.append(line)
        if slide_texts:
            parts.append(f"=== Slide {slide_num} ===")
            parts.extend(slide_texts)
    return "\n".join(parts) if parts else ""


def _read_csv(path: Path) -> str:
    """Read a CSV file as tab-separated text."""
    rows: list[str] = []
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for row in reader:
            line = "\t".join(cell.strip() for cell in row)
            if line:
                rows.append(line)
    return "\n".join(rows) if rows else ""


def _read_pdf(path: Path, use_ocr: bool = False) -> str:
    """Read PDF text. Uses OCR if use_ocr is True or if PDF is scanned."""
    if use_ocr:
        ocr_text = _ocr_pdf(path)
        if ocr_text.strip():
            return ocr_text

    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("PDF parsing requires optional dependency 'pypdf'. Install with: pip install -e .[pdf]") from exc

    try:
        reader = PdfReader(str(path))
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as exc:
        raise RuntimeError(f"Failed to read PDF {path.name}: {exc}") from exc

    # If text is too short (scanned PDF), fall back to OCR if available
    if len(text.strip()) < 50:
        try:
            ocr_text = _ocr_pdf(path)
            if ocr_text.strip():
                return ocr_text
        except Exception:
            pass

    return text


def _read_image_with_ocr(path: Path) -> str:
    """Read text from an image using OCR."""
    try:
        from rbs_rag.ocr.service import get_ocr_service
        service = get_ocr_service()
        result = service.process(path)
        return result.full_text if result and result.full_text else ""
    except Exception as exc:
        return f"[OCR Error: {exc}]"


def _read_xlsx(path: Path) -> str:
    """Read all cell values from an .xlsx workbook."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        rows = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows.append(f"=== Sheet: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = "\t".join(cells).strip()
                if line:
                    rows.append(line)
        wb.close()
        return "\n".join(rows) if rows else ""
    except Exception as exc:
        return f"[Error reading .xlsx: {exc}]"


def _read_xls(path: Path) -> str:
    """Read all cell values from a legacy .xls workbook."""
    try:
        import xlrd
        wb = xlrd.open_workbook(str(path))
        rows = []
        for sheet_idx in range(wb.nsheets):
            ws = wb.sheet_by_index(sheet_idx)
            rows.append(f"=== Sheet: {ws.name} ===")
            for row_idx in range(ws.nrows):
                cells = [str(ws.cell_value(row_idx, c)) for c in range(ws.ncols)]
                line = "\t".join(cells).strip()
                if line:
                    rows.append(line)
        return "\n".join(rows) if rows else ""
    except Exception as exc:
        return f"[Error reading .xls: {exc}]"


def _ocr_pdf(path: Path) -> str:
    """Run OCR on a PDF and return extracted text."""
    try:
        from rbs_rag.ocr.service import get_ocr_service
        service = get_ocr_service()
        result = service.process(path)
        return result.full_text if result and result.full_text else ""
    except Exception as exc:
        return f"[OCR Error: {exc}]"
